import base64
import io
import json
import logging
import re
import zipfile
from datetime import UTC, datetime, timedelta
from typing import Any, Optional

import httpx
from langchain.chat_models import init_chat_model
from langchain_core.messages import HumanMessage, SystemMessage
from PIL import Image, ImageDraw

from app.core.config import get_settings
from app.domain.external.file import FileStorage
from app.domain.models.agent import Agent
from app.domain.models.event import MessageEvent, TitleEvent
from app.domain.models.file import FileInfo
from app.domain.repositories.session_repository import SessionRepository
from app.infrastructure.models.documents import ScheduleDocument

logger = logging.getLogger(__name__)


class ModeOutputService:
    def __init__(
        self,
        agent: Agent,
        session_id: str,
        user_id: str,
        file_storage: FileStorage,
        session_repository: SessionRepository,
    ):
        self._agent = agent
        self._session_id = session_id
        self._user_id = user_id
        self._file_storage = file_storage
        self._session_repository = session_repository
        self._settings = get_settings()
        self._model = self._build_model()

    def _build_model(self):
        kwargs: dict[str, Any] = dict(
            model=self._agent.model_name,
            model_provider=self._settings.model_provider,
            temperature=self._agent.temperature,
            max_tokens=self._agent.max_tokens,
            base_url=self._settings.api_base,
        )
        if self._settings.extra_headers:
            kwargs["default_headers"] = self._settings.extra_headers
        return init_chat_model(**kwargs)

    async def _invoke(self, system_prompt: str, user_prompt: str) -> str:
        response = await self._model.ainvoke(
            [SystemMessage(content=system_prompt), HumanMessage(content=user_prompt)]
        )
        content = getattr(response, "content", "")
        if isinstance(content, list):
            return "\n".join(str(item) for item in content)
        return str(content or "")

    def _parse_json_block(self, text: str, fallback: Any) -> Any:
        try:
            return json.loads(text)
        except Exception:
            pass
        match = re.search(r"```json\s*(.*?)```", text, re.DOTALL)
        if match:
            try:
                return json.loads(match.group(1))
            except Exception:
                pass
        return fallback

    async def _upload_session_file(
        self,
        filename: str,
        content: bytes,
        content_type: str,
        metadata: Optional[dict[str, Any]] = None,
    ) -> FileInfo:
        file_info = await self._file_storage.upload_file(
            io.BytesIO(content),
            filename,
            self._user_id,
            content_type=content_type,
            metadata=metadata or {},
        )
        await self._session_repository.add_file(self._session_id, file_info)
        return file_info

    def _title_from_prompt(self, prefix: str, prompt: str) -> str:
        cleaned = re.sub(r"\s+", " ", prompt).strip()
        if not cleaned:
            return prefix
        return f"{prefix}: {cleaned[:56]}".strip()

    async def run_slides_mode(self, prompt: str, slides_template: str = "Minimal"):
        from pptx import Presentation
        from pptx.util import Inches

        yield TitleEvent(title=self._title_from_prompt("Slides", prompt))

        raw_outline = await self._invoke(
            "Return JSON only with a title and 4-6 slides. Each slide needs title and 2-4 bullet strings.",
            json.dumps({"prompt": prompt, "style": slides_template}),
        )
        outline = self._parse_json_block(raw_outline, {})
        slides = outline.get("slides") if isinstance(outline, dict) else None
        if not isinstance(slides, list) or not slides:
            slides = [
                {"title": "Overview", "bullets": [prompt, "Why it matters", "Recommended direction"]},
                {"title": "Audience", "bullets": ["Primary audience", "Key needs", "Buying triggers"]},
                {"title": "Solution", "bullets": ["Core workflow", "Differentiators", "Expected impact"]},
                {"title": "Next steps", "bullets": ["Launch plan", "Execution checklist", "Success metrics"]},
            ]

        presentation = Presentation()
        for slide_data in slides[:6]:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = str(slide_data.get("title") or "Slide")
            body = slide.placeholders[1].text_frame
            body.clear()
            for idx, bullet in enumerate(slide_data.get("bullets", [])[:4]):
                paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
                paragraph.text = str(bullet)
            # Add a subtle style accent block so templates feel different.
            accent = slide.shapes.add_textbox(Inches(9), Inches(0.2), Inches(1), Inches(0.4))
            accent.text_frame.text = slides_template

        output = io.BytesIO()
        presentation.save(output)
        pptx_file = await self._upload_session_file(
            f"{self._slug(prompt)}_deck.pptx",
            output.getvalue(),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            metadata={"kind": "slides", "template": slides_template},
        )

        summary = f"Built a {slides_template} presentation deck for `{prompt}`."
        yield MessageEvent(message=summary, attachments=[pptx_file])

    async def run_website_mode(self, prompt: str, website_category: str = "SaaS"):
        yield TitleEvent(title=self._title_from_prompt("Website", prompt))

        raw_site = await self._invoke(
            "Return JSON only with keys title, hero, features, cta, html, css, js. Keep the site premium and concise.",
            json.dumps({"prompt": prompt, "category": website_category}),
        )
        site = self._parse_json_block(raw_site, {})
        title = site.get("title") if isinstance(site, dict) else None
        hero = site.get("hero") if isinstance(site, dict) else None
        features = site.get("features") if isinstance(site, dict) else None
        cta = site.get("cta") if isinstance(site, dict) else None

        if not isinstance(features, list) or not features:
            features = ["Intentional monochrome design", "AI-assisted workflows", "Fast shipping-ready pages"]

        css = str(site.get("css") or self._default_website_css())
        js = str(site.get("js") or self._default_website_js())
        html = str(
            site.get("html")
            or self._default_website_html(
                title or "Forge Website Preview",
                hero or prompt,
                [str(item) for item in features[:4]],
                cta or "Launch with Forge",
                css,
                js,
            )
        )

        slug = self._slug(prompt)
        html_file = await self._upload_session_file(f"{slug}_preview.html", html.encode("utf-8"), "text/html", metadata={"kind": "website_preview"})
        css_file = await self._upload_session_file(f"{slug}_styles.css", css.encode("utf-8"), "text/css", metadata={"kind": "website_asset"})
        js_file = await self._upload_session_file(f"{slug}_script.js", js.encode("utf-8"), "application/javascript", metadata={"kind": "website_asset"})

        archive_buffer = io.BytesIO()
        with zipfile.ZipFile(archive_buffer, "w", zipfile.ZIP_DEFLATED) as archive:
            archive.writestr("index.html", html)
            archive.writestr("styles.css", css)
            archive.writestr("script.js", js)
        zip_file = await self._upload_session_file(f"{slug}_site.zip", archive_buffer.getvalue(), "application/zip", metadata={"kind": "website_bundle"})

        summary = (
            f"Built a {website_category} website starter for `{prompt}`. "
            "Open the attached HTML file for the live preview, or download the ZIP for the full bundle."
        )
        yield MessageEvent(message=summary, attachments=[html_file, css_file, js_file, zip_file])

    async def run_design_mode(self, prompt: str, design_model: str = "Forge Image v1"):
        yield TitleEvent(title=self._title_from_prompt("Design", prompt))

        image_bytes = await self._try_nvidia_image(prompt, design_model)
        filename = f"{self._slug(prompt)}_{'v2' if 'v2' in design_model.lower() else 'v1'}.png"
        content_type = "image/png"

        if image_bytes is None:
            image_bytes = self._fallback_design_png(prompt, design_model)

        image_file = await self._upload_session_file(
            filename,
            image_bytes,
            content_type,
            metadata={"kind": "design_image", "design_model": design_model},
        )
        yield MessageEvent(message=f"Generated a design concept with {design_model} for `{prompt}`.", attachments=[image_file])

    async def run_spreadsheet_mode(self, prompt: str):
        from openpyxl import Workbook

        yield TitleEvent(title=self._title_from_prompt("Spreadsheet", prompt))
        raw = await self._invoke(
            "Return JSON only with workbook_title, sheet_name, columns, and rows. Keep rows concise and useful.",
            json.dumps({"prompt": prompt}),
        )
        parsed = self._parse_json_block(raw, {})
        columns = parsed.get("columns") if isinstance(parsed, dict) else None
        rows = parsed.get("rows") if isinstance(parsed, dict) else None
        if not isinstance(columns, list) or not columns:
            columns = ["Category", "Owner", "Status", "Notes"]
        if not isinstance(rows, list) or not rows:
            rows = [
                ["Planning", "Forge", "Active", "Starter row generated by Forge"],
                ["Execution", "Ops", "Queued", "Adapt this workbook to your exact data"],
            ]

        workbook = Workbook()
        sheet = workbook.active
        sheet.title = str(parsed.get("sheet_name") or "Forge Sheet")
        sheet.append([str(col) for col in columns])
        for row in rows[:50]:
            if isinstance(row, list):
                sheet.append([str(cell) for cell in row[: len(columns)]])

        workbook_bytes = io.BytesIO()
        workbook.save(workbook_bytes)
        xlsx = await self._upload_session_file(
            f"{self._slug(prompt)}.xlsx",
            workbook_bytes.getvalue(),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            metadata={"kind": "spreadsheet"},
        )
        summary = await self._upload_session_file(
            f"{self._slug(prompt)}_summary.md",
            self._markdown_export("Spreadsheet output", prompt, raw).encode("utf-8"),
            "text/markdown",
            metadata={"kind": "spreadsheet_notes"},
        )
        yield MessageEvent(message=f"Built a spreadsheet output for `{prompt}`.", attachments=[xlsx, summary])

    async def run_video_mode(self, prompt: str):
        yield TitleEvent(title=self._title_from_prompt("Video", prompt))
        content = await self._invoke(
            "Write a structured video concept with overview, scene-by-scene breakdown, timestamps, and CTA.",
            prompt,
        )
        export = await self._upload_session_file(
            f"{self._slug(prompt)}_video_brief.md",
            self._markdown_export("Video brief", prompt, content).encode("utf-8"),
            "text/markdown",
            metadata={"kind": "video_brief"},
        )
        yield MessageEvent(message=content, attachments=[export])

    async def run_audio_mode(self, prompt: str):
        yield TitleEvent(title=self._title_from_prompt("Audio", prompt))
        content = await self._invoke(
            "Write polished audio output with sections, timing guidance, and delivery notes where useful.",
            prompt,
        )
        export = await self._upload_session_file(
            f"{self._slug(prompt)}_audio_brief.md",
            self._markdown_export("Audio brief", prompt, content).encode("utf-8"),
            "text/markdown",
            metadata={"kind": "audio_brief"},
        )
        yield MessageEvent(message=content, attachments=[export])

    async def run_visualization_mode(self, prompt: str):
        yield TitleEvent(title=self._title_from_prompt("Visualization", prompt))
        content = await self._invoke(
            "Return JSON only with title, summary, labels, and series. Labels should be short strings, series numeric.",
            prompt,
        )
        parsed = self._parse_json_block(content, {})
        labels = parsed.get("labels") if isinstance(parsed, dict) else None
        series = parsed.get("series") if isinstance(parsed, dict) else None
        if not isinstance(labels, list) or not labels:
            labels = ["Jan", "Feb", "Mar", "Apr", "May"]
        if not isinstance(series, list) or not series:
            series = [14, 22, 18, 31, 27]
        title = str(parsed.get("title") or "Forge Visualization")
        html = self._build_chart_html(title, [str(label) for label in labels], [float(value) for value in series])
        export = await self._upload_session_file(
            f"{self._slug(prompt)}_visualization.html",
            html.encode("utf-8"),
            "text/html",
            metadata={"kind": "visualization"},
        )
        yield MessageEvent(message=f"Built a visualization for `{prompt}`.", attachments=[export])

    async def run_develop_apps_mode(self, prompt: str):
        yield TitleEvent(title=self._title_from_prompt("Develop App", prompt))
        spec = await self._invoke(
            "Return JSON only with app_name, summary, files [{path, content}], and launch_notes.",
            prompt,
        )
        parsed = self._parse_json_block(spec, {})
        app_name = str(parsed.get("app_name") or "forge_app")
        files = parsed.get("files") if isinstance(parsed, dict) else None
        if not isinstance(files, list) or not files:
            files = [
                {"path": "README.md", "content": f"# {app_name}\n\nGenerated by Forge.\n\n{prompt}"},
                {"path": "src/index.js", "content": "console.log('Forge starter app');\n"},
            ]

        archive_buffer = io.BytesIO()
        with zipfile.ZipFile(archive_buffer, "w", zipfile.ZIP_DEFLATED) as archive:
            for file in files[:20]:
                archive.writestr(str(file.get("path") or "README.md"), str(file.get("content") or ""))
        zip_file = await self._upload_session_file(
            f"{self._slug(app_name)}_starter.zip",
            archive_buffer.getvalue(),
            "application/zip",
            metadata={"kind": "starter_app"},
        )
        summary_file = await self._upload_session_file(
            f"{self._slug(app_name)}_spec.md",
            self._markdown_export("App spec", prompt, spec).encode("utf-8"),
            "text/markdown",
            metadata={"kind": "app_spec"},
        )
        yield MessageEvent(message=str(parsed.get("summary") or spec), attachments=[zip_file, summary_file])

    async def run_playbook_mode(self, prompt: str):
        yield TitleEvent(title=self._title_from_prompt("Playbook", prompt))
        content = await self._invoke(
            "Write a detailed operational playbook with objective, triggers, roles, steps, and QA checklist.",
            prompt,
        )
        export = await self._upload_session_file(
            f"{self._slug(prompt)}_playbook.md",
            self._markdown_export("Playbook", prompt, content).encode("utf-8"),
            "text/markdown",
            metadata={"kind": "playbook"},
        )
        yield MessageEvent(message=content, attachments=[export])

    async def run_schedule_mode(self, prompt: str):
        yield TitleEvent(title=self._title_from_prompt("Schedule", prompt))
        cron_expression, next_run = self._parse_schedule(prompt)
        schedule_doc = ScheduleDocument(
            schedule_id=self._slug(f"{prompt}-{datetime.now(UTC).isoformat()}")[:16],
            user_id=self._user_id,
            name=(prompt[:60] or "New schedule").strip(),
            description=prompt,
            schedule_text=prompt,
            cron_expression=cron_expression,
            next_run=next_run,
            status="active",
        )
        await schedule_doc.insert()
        export = await self._upload_session_file(
            f"{self._slug(prompt)}_schedule.md",
            self._markdown_export("Scheduled task", prompt, f"Cron: {cron_expression}\nNext run: {next_run.isoformat()}").encode("utf-8"),
            "text/markdown",
            metadata={"kind": "schedule"},
        )
        yield MessageEvent(
            message=f"Saved a scheduled task for `{prompt}`.\n\nCron: `{cron_expression}`\nNext run: `{next_run.isoformat()}`",
            attachments=[export],
        )

    def _markdown_export(self, title: str, prompt: str, body: str) -> str:
        generated_at = datetime.now(UTC).strftime("%Y-%m-%d %H:%M UTC")
        return f"# {title}\n\nGenerated by Forge on {generated_at}\n\n## Prompt\n\n{prompt}\n\n## Output\n\n{body}\n"

    def _parse_schedule(self, schedule_text: str) -> tuple[str, datetime]:
        text = schedule_text.lower()
        now = datetime.now(UTC)
        hour = 9
        minute = 0
        match = re.search(r"(\d{1,2})(?::(\d{2}))?\s*(am|pm)?", text)
        if match:
            hour = int(match.group(1))
            minute = int(match.group(2) or 0)
            meridiem = match.group(3)
            if meridiem == "pm" and hour < 12:
                hour += 12
            if meridiem == "am" and hour == 12:
                hour = 0
        next_run = now.replace(hour=hour, minute=minute, second=0, microsecond=0)
        if "weekly" in text or "monday" in text:
            next_run = next_run.replace(hour=hour, minute=minute) + timedelta(days=7)
            return f"{minute} {hour} * * 1", next_run
        if next_run <= now:
            next_run = next_run.replace(hour=hour, minute=minute) + timedelta(days=1)
        return f"{minute} {hour} * * *", next_run

    def _build_chart_html(self, title: str, labels: list[str], series: list[float]) -> str:
        labels_json = json.dumps(labels)
        series_json = json.dumps(series)
        return f"""<!doctype html>
<html lang="en">
  <head>
    <meta charset="utf-8" />
    <meta name="viewport" content="width=device-width, initial-scale=1" />
    <title>{title}</title>
    <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
    <style>
      body {{ margin: 0; font-family: Arial, sans-serif; background: #faf9f7; color: #111827; }}
      .shell {{ max-width: 960px; margin: 0 auto; padding: 48px 24px; }}
      .card {{ background: white; border: 1px solid rgba(0,0,0,.08); border-radius: 24px; padding: 24px; box-shadow: 0 14px 36px rgba(0,0,0,.05); }}
      h1 {{ margin: 0 0 24px; font: 400 40px Georgia, 'Times New Roman', serif; }}
    </style>
  </head>
  <body>
    <main class="shell">
      <div class="card">
        <h1>{title}</h1>
        <canvas id="chart"></canvas>
      </div>
    </main>
    <script>
      new Chart(document.getElementById('chart'), {{
        type: 'bar',
        data: {{
          labels: {labels_json},
          datasets: [{{ label: {json.dumps(title)}, data: {series_json}, backgroundColor: '#111827', borderRadius: 10 }}]
        }},
        options: {{ responsive: true, plugins: {{ legend: {{ display: false }} }} }}
      }});
    </script>
  </body>
</html>"""

    async def _try_nvidia_image(self, prompt: str, design_model: str) -> Optional[bytes]:
        if not self._settings.api_key or not self._settings.api_base:
            return None
        image_model = "nvidia/cosmos-predict2" if "v2" in design_model.lower() else "black-forest-labs/flux.1-dev"
        url = f"{self._settings.api_base.rstrip('/')}/images/generations"
        payload = {"model": image_model, "prompt": prompt, "size": "1024x1024"}
        headers = {"Authorization": f"Bearer {self._settings.api_key}"}
        try:
            async with httpx.AsyncClient(timeout=20.0) as client:
                response = await client.post(url, json=payload, headers=headers)
                if response.status_code >= 400:
                    logger.warning("NVIDIA image generation failed: %s", response.text[:240])
                    return None
                data = response.json()
                image_data = (data.get("data") or [{}])[0]
                if image_data.get("b64_json"):
                    return base64.b64decode(image_data["b64_json"])
                if image_data.get("url"):
                    download = await client.get(image_data["url"])
                    if download.status_code < 400:
                        return download.content
        except Exception as exc:
            logger.warning("Image generation fallback triggered: %s", exc)
        return None

    def _fallback_design_png(self, prompt: str, design_model: str) -> bytes:
        image = Image.new("RGB", (1200, 800), color=(17, 24, 39))
        draw = ImageDraw.Draw(image)
        draw.rectangle((60, 60, 1140, 740), outline=(229, 231, 235), width=3)
        draw.text((100, 110), "Forge Design Preview", fill=(255, 255, 255))
        draw.text((100, 180), design_model, fill=(147, 197, 253))
        wrapped = re.sub(r"(.{1,48})(\s+|$)", "\\1\n", prompt).strip()
        draw.text((100, 260), wrapped[:420], fill=(229, 231, 235))
        buffer = io.BytesIO()
        image.save(buffer, format="PNG")
        return buffer.getvalue()

    def _slug(self, prompt: str) -> str:
        return re.sub(r"[^a-z0-9]+", "_", prompt.lower()).strip("_")[:32] or "forge_output"

    def _default_website_css(self) -> str:
        return """
:root {
  --bg: #0b0b0c;
  --panel: #151517;
  --text: #fafaf9;
  --muted: #b8b8bd;
  --line: rgba(255,255,255,0.08);
}
* { box-sizing: border-box; }
body {
  margin: 0;
  font-family: Georgia, 'Times New Roman', serif;
  background: radial-gradient(circle at top, #1f2937, #0b0b0c 58%);
  color: var(--text);
}
.shell { max-width: 1100px; margin: 0 auto; padding: 72px 24px; }
.hero { padding: 56px; border: 1px solid var(--line); border-radius: 32px; background: rgba(255,255,255,0.03); }
.eyebrow { font: 600 12px/1.2 Arial, sans-serif; letter-spacing: .18em; text-transform: uppercase; color: var(--muted); }
h1 { font-size: clamp(44px, 7vw, 84px); line-height: .95; margin: 18px 0; }
p { font: 400 18px/1.7 Arial, sans-serif; color: var(--muted); max-width: 720px; }
.grid { display: grid; grid-template-columns: repeat(auto-fit, minmax(220px, 1fr)); gap: 18px; margin-top: 36px; }
.card { border: 1px solid var(--line); border-radius: 24px; padding: 22px; background: rgba(255,255,255,0.03); }
.cta { margin-top: 28px; display: inline-flex; padding: 14px 20px; border-radius: 999px; background: white; color: #111827; font: 600 14px/1 Arial, sans-serif; text-decoration: none; }
        """.strip()

    def _default_website_js(self) -> str:
        return """
document.querySelectorAll('[data-animate]').forEach((node, index) => {
  node.animate([{ opacity: 0, transform: 'translateY(16px)' }, { opacity: 1, transform: 'translateY(0)' }], {
    duration: 500,
    delay: index * 100,
    fill: 'forwards',
    easing: 'ease-out'
  });
});
        """.strip()

    def _default_website_html(self, title: str, hero: str, features: list[str], cta: str, css: str, js: str) -> str:
        feature_markup = "\n".join(f"<div class='card' data-animate><strong>{item}</strong></div>" for item in features)
        return f"""<!doctype html>
<html lang="en">
  <head>
    <meta charset="utf-8" />
    <meta name="viewport" content="width=device-width, initial-scale=1" />
    <title>{title}</title>
    <style>{css}</style>
  </head>
  <body>
    <main class="shell">
      <section class="hero" data-animate>
        <div class="eyebrow">Forge Generated Site</div>
        <h1>{title}</h1>
        <p>{hero}</p>
        <a class="cta" href="#features">{cta}</a>
      </section>
      <section id="features" class="grid">
        {feature_markup}
      </section>
    </main>
    <script>{js}</script>
  </body>
</html>"""
